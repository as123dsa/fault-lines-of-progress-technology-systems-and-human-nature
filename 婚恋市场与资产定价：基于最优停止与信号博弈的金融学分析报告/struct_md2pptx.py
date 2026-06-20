#!/usr/bin/env python3
"""
将结构化 Markdown 报告转换为 PowerPoint 演示文稿，支持指定模板文件。
格式要求：
- 页面分隔符：---
- 页面标题：### PPT N：标题内容
- 字段：**标题：**、**副标题：**、**汇报人/机构：**、**内容：**、**备注/解说词：**
- 内容区支持无序列表（- 开头）和普通段落
- 以 "--- **（PPT结束）**" 结尾的页面将被忽略
"""

import argparse
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def parse_md_blocks(md_text):
    """
    按 '---' 分割 Markdown，返回每个块的原始文本列表。
    同时过滤掉结尾标记为 "**（PPT结束）**" 的块。
    """
    # 使用正则分割，但不要分割开头的可选 ---
    # 分割符为单独一行 '---'（前后可能有空白）
    blocks = re.split(r'\n---\n', md_text.strip())
    # 如果第一个块以 --- 开头（文件头），去掉它（但本文件第一个块就是第一页）
    if blocks and blocks[0].startswith('---'):
        blocks[0] = blocks[0][3:].strip()
    filtered_blocks = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        # 检查是否是结束标记块：内容仅为 "**（PPT结束）**" 或类似
        if re.match(r'^\*\*（PPT结束）\*\*$', block):
            continue
        # 如果整个块只有 "---" 和结束标记，但已经被分割，忽略
        filtered_blocks.append(block)
    return filtered_blocks


def parse_field(block, field_name):
    """
    从块中提取某个字段（如 **标题：**）的内容。
    返回从该字段标记后到下一个字段标记（或块结束）之间的文本，去除首尾空白。
    支持跨行。
    """
    pattern = rf'^\*\*{re.escape(field_name)}：\*\*\s*(.*?)(?=^\*\*[^*]+：\*\*|\Z)'
    match = re.search(pattern, block, re.DOTALL | re.MULTILINE)
    if match:
        return match.group(1).strip()
    return None


def parse_list_items(content_text):
    """
    将内容文本解析为段落列表，每个元素可以是 (text, bullet_flag)。
    支持无序列表（以 "- " 开头）和普通文本行。
    """
    lines = content_text.splitlines()
    items = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # 检查是否是无序列表项
        if line.startswith('- '):
            items.append((line[2:], True))
        else:
            items.append((line, False))
    return items


def add_paragraph_with_format(tf, text, bullet=False, level=0, font_size=18, bold=False, italic=False):
    """向文本框添加段落，支持内联加粗（**text**）"""
    p = tf.add_paragraph()
    # 简单处理内联加粗
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.size = Pt(font_size)
        else:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            if bold:
                run.font.bold = True
            if italic:
                run.font.italic = True
    p.bullet = bullet
    p.level = level
    return p


def set_slide_notes(slide, notes_text):
    """添加备注到幻灯片"""
    if not notes_text:
        return
    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = notes_text
    except Exception as e:
        print(f"警告：添加备注失败 - {e}")


def create_cover_slide(prs, title_text, subtitle_text, author_text):
    """
    创建封面幻灯片。
    使用模板中的第一个布局（通常索引0为标题幻灯片）。
    如果模板中不存在标题或副标题占位符，则回退到自定义文本框。
    """
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    # 设置标题
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    else:
        # 如果没有标题占位符，手动添加文本框
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tb.text_frame.paragraphs[0].text = title_text
        tb.text_frame.paragraphs[0].font.size = Pt(44)

    # 设置副标题区域（可能是占位符索引1或直接查找）
    subtitle_placeholder = None
    if len(slide.placeholders) > 1:
        subtitle_placeholder = slide.placeholders[1]
    if subtitle_placeholder:
        combined = subtitle_text
        if author_text:
            combined = f"{subtitle_text}\n{author_text}" if subtitle_text else author_text
        subtitle_placeholder.text = combined
    else:
        # 手动添加副文本框
        if subtitle_text or author_text:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
            tf = tb.text_frame
            tf.paragraphs[0].text = subtitle_text
            if author_text:
                p = tf.add_paragraph()
                p.text = author_text
                p.font.size = Pt(18)
    return slide


def add_content_slide(prs, slide_title, content_items, notes_text):
    """创建普通内容页，使用模板中的“标题和内容”布局（索引1）"""
    try:
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    except IndexError:
        # 如果模板没有布局1，使用空白布局手动添加
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        # 手动加标题
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tb.text_frame.paragraphs[0].text = slide_title
        tb.text_frame.paragraphs[0].font.size = Pt(32)
        # 手动加内容区域
        content_tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = content_tb.text_frame
        tf.word_wrap = True
        for text, bullet in content_items:
            add_paragraph_with_format(tf, text, bullet=bullet)
        set_slide_notes(slide, notes_text)
        return slide

    slide = prs.slides.add_slide(slide_layout)
    # 设置标题
    if slide.shapes.title:
        slide.shapes.title.text = slide_title
    # 清空内容占位符并重新填充（避免格式混乱）
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # 内容占位符通常索引为1
            content_placeholder = shape
            break
    if content_placeholder:
        # 清除原有段落
        tf = content_placeholder.text_frame
        tf.clear()
        for text, bullet in content_items:
            add_paragraph_with_format(tf, text, bullet=bullet)
    else:
        # 如果没有内容占位符，手动添加文本框
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        for text, bullet in content_items:
            add_paragraph_with_format(tf, text, bullet=bullet)

    set_slide_notes(slide, notes_text)
    return slide


def md_to_pptx(md_path, pptx_path, template_path=None):
    # 读取 Markdown
    with open(md_path, 'r', encoding='utf-8') as f:
        md_text = f.read()

    blocks = parse_md_blocks(md_text)

    # 加载模板
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"使用模板：{template_path}")
    else:
        prs = Presentation()
        if template_path:
            print(f"警告：模板文件 {template_path} 不存在，使用默认空白模板。")
        else:
            print("未提供模板，使用默认空白模板。")
    # 可选：设置默认页面尺寸（如果模板未定义，保持模板原有尺寸）
    # prs.slide_width = Inches(10)  # 建议不覆盖模板尺寸
    # prs.slide_height = Inches(7.5)

    # 遍历每个页面块
    for idx, block in enumerate(blocks):
        # 提取页面标题
        lines = block.split('\n')
        title_line = lines[0].strip() if lines else ""
        page_title = ""
        if title_line.startswith('### PPT'):
            match = re.search(r'：\s*(.*)', title_line)
            if match:
                page_title = match.group(1).strip()
            else:
                page_title = title_line[3:].strip()
        else:
            page_title = "未命名"

        # 解析各个字段
        title_field = parse_field(block, '标题')
        subtitle_field = parse_field(block, '副标题')
        author_field = parse_field(block, '汇报人/机构')
        content_field = parse_field(block, '内容')
        notes_field = parse_field(block, '备注/解说词')

        # 判断是否为封面（第一个页面且包含标题、副标题或作者字段）
        is_cover = (idx == 0 and (title_field or subtitle_field or author_field))

        if is_cover:
            slide_title = title_field or page_title
            slide = create_cover_slide(prs, slide_title, subtitle_field or "", author_field or "")
            set_slide_notes(slide, notes_field)
        else:
            # 普通内容页
            content_items = []
            if content_field:
                content_items = parse_list_items(content_field)
            # 如果没有内容字段，尝试使用块中除标题和元字段之外的行
            if not content_items:
                temp_text = block
                for f in ['标题', '副标题', '汇报人/机构', '内容', '备注/解说词']:
                    temp_text = re.sub(rf'^\*\*{f}：\*\*.*?(?=^\*\*[^*]+：\*\*|\Z)', '', temp_text, flags=re.DOTALL | re.MULTILINE)
                for line in temp_text.split('\n'):
                    line = line.strip()
                    if line and not line.startswith('###'):
                        content_items.append((line, False))
            slide = add_content_slide(prs, page_title, content_items, notes_field)

    # 保存文件
    prs.save(pptx_path)
    print(f"成功生成 {pptx_path}")


def main():
    parser = argparse.ArgumentParser(description="将结构化 Markdown 报告转换为 PowerPoint，支持模板")
    parser.add_argument("input", help="输入的 Markdown 文件路径")
    parser.add_argument("output", help="输出的 PowerPoint 文件路径 (.pptx)")
    parser.add_argument("-t", "--template", help="PowerPoint 模板文件路径（.pptx），将沿用其中的主题和母版")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"错误：输入文件 {args.input} 不存在")
        return 1
    try:
        md_to_pptx(args.input, args.output, args.template)
    except Exception as e:
        print(f"转换失败：{e}")
        import traceback
        traceback.print_exc()
        return 1
    return 0


if __name__ == "__main__":
    exit(main())